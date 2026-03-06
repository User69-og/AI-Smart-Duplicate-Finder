import os
import re
import hashlib
import datetime
import csv
from collections import defaultdict

from PIL import Image
import imagehash
import PyPDF2
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity


TEXT_EXTENSIONS  = {".txt", ".docx", ".xlsx", ".pptx", ".pdf"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".gif", ".webp"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".flac", ".ogg", ".m4a"}


# ============================================================
# SHA-256 EXACT DUPLICATE HASH
# ============================================================

def sha256_hash(path):
    h = hashlib.sha256()
    try:
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return None


# ============================================================
# TEXT CLEANING
# ============================================================

def clean_text(text):
    text = text.lower()
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[^a-z0-9 ]", " ", text)
    return text.strip()


# ============================================================
# TEXT EXTRACTION
# ============================================================

def extract_text(path):
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            doc = Document(path)
            return clean_text(" ".join(p.text for p in doc.paragraphs))

        if ext == ".xlsx":
            wb = load_workbook(path, read_only=True, data_only=True)
            text = []
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell:
                            text.append(str(cell))
            return clean_text(" ".join(text))

        if ext == ".pptx":
            prs = Presentation(path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return clean_text(" ".join(text))

        if ext == ".pdf":
            reader = PyPDF2.PdfReader(path)
            text = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text.append(t)
            return clean_text(" ".join(text))

        if ext == ".txt":
            with open(path, "r", encoding="utf8", errors="ignore") as f:
                return clean_text(f.read())

    except Exception as e:
        print(f"[WARN] extract_text failed for {path}: {e}")
        return ""
    return ""


# ============================================================
# IMAGE PHASH
# ============================================================

def image_hash(path):
    try:
        return imagehash.phash(Image.open(path))
    except Exception:
        return None


def get_image_thumbnail(path, size=(120, 120)):
    try:
        img = Image.open(path)
        img.thumbnail(size, Image.LANCZOS)
        return img
    except Exception:
        return None


# ============================================================
# AUDIO FINGERPRINTING
# ============================================================

def audio_fingerprint(path):
    try:
        import librosa
        import numpy as np
        y, sr = librosa.load(path, sr=None, mono=True, duration=30)
        mfcc = librosa.feature.mfcc(y=y, sr=sr, n_mfcc=20)
        return mfcc.mean(axis=1)
    except Exception:
        return None


def audio_similarity(fp1, fp2):
    try:
        import numpy as np
        n1 = fp1 / (np.linalg.norm(fp1) + 1e-10)
        n2 = fp2 / (np.linalg.norm(fp2) + 1e-10)
        return float(np.dot(n1, n2))
    except Exception:
        return 0.0


# ============================================================
# FAISS VECTOR INDEX
# ============================================================

def build_faiss_index(embeddings):
    import faiss
    import numpy as np
    embeddings = np.array(embeddings).astype("float32")
    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings)
    return index


# ============================================================
# GRAPH CLUSTERING
# ============================================================

def cluster_pairs(pairs):
    graph = defaultdict(set)
    for a, b in pairs:
        graph[a].add(b)
        graph[b].add(a)

    visited  = set()
    clusters = []
    for node in graph:
        if node not in visited:
            stack   = [node]
            cluster = []
            while stack:
                n = stack.pop()
                if n not in visited:
                    visited.add(n)
                    cluster.append(n)
                    stack.extend(graph[n] - visited)
            clusters.append(cluster)
    return clusters


# ============================================================
# RECOMMENDATION ENGINE
# ============================================================

def score_file(path):
    score   = 0
    reasons = []
    try:
        score += os.path.getsize(path)
        reasons.append("larger file size")
        score += os.path.getmtime(path)
        reasons.append("more recent modification")
        name = os.path.basename(path).lower()
        if "final" in name:
            score += 10_000_000_000
            reasons.append("marked as final")
        if "copy" in name or "backup" in name:
            score -= 5_000_000_000
            reasons.append("likely backup copy")
    except Exception:
        pass
    return score, reasons


def recommend_cluster(cluster):
    scored = []
    for f in cluster:
        s, r = score_file(f)
        scored.append((f, s, r))
    scored.sort(key=lambda x: x[1], reverse=True)
    return scored[0], scored[1:]


# ============================================================
# FILE METADATA
# ============================================================

def get_file_metadata(path):
    try:
        stat       = os.stat(path)
        size_bytes = stat.st_size
        modified   = datetime.datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M")
        try:
            created = datetime.datetime.fromtimestamp(stat.st_ctime).strftime("%Y-%m-%d %H:%M")
        except Exception:
            created = "—"
        ext = os.path.splitext(path)[1].lower()

        if size_bytes < 1024:
            size_str = f"{size_bytes} B"
        elif size_bytes < 1024 ** 2:
            size_str = f"{size_bytes / 1024:.1f} KB"
        elif size_bytes < 1024 ** 3:
            size_str = f"{size_bytes / 1024 ** 2:.1f} MB"
        else:
            size_str = f"{size_bytes / 1024 ** 3:.2f} GB"

        return {
            "size_bytes": size_bytes,
            "size_str":   size_str,
            "modified":   modified,
            "created":    created,
            "ext":        ext,
        }
    except Exception:
        return {"size_bytes": 0, "size_str": "?", "modified": "?", "created": "?", "ext": "?"}


# ============================================================
# SCAN SUMMARY
# ============================================================

def compute_summary(exact_groups, near_clusters, all_files):
    total_files  = len(all_files)
    wasted_bytes = 0
    dup_files    = set()

    for group in exact_groups:
        for f in group[1:]:
            try:
                wasted_bytes += os.path.getsize(f)
            except Exception:
                pass
            dup_files.add(f)

    for cluster in near_clusters:
        _, remove = recommend_cluster(cluster)
        for f, _, _ in remove:
            dup_files.add(f)
            try:
                wasted_bytes += os.path.getsize(f)
            except Exception:
                pass

    if wasted_bytes < 1024 ** 2:
        wasted_str = f"{wasted_bytes / 1024:.1f} KB"
    elif wasted_bytes < 1024 ** 3:
        wasted_str = f"{wasted_bytes / 1024 ** 2:.1f} MB"
    else:
        wasted_str = f"{wasted_bytes / 1024 ** 3:.2f} GB"

    return {
        "total_files":         total_files,
        "duplicate_files":     len(dup_files),
        "wasted_bytes":        wasted_bytes,
        "wasted_str":          wasted_str,
        "exact_groups_count":  len(exact_groups),
        "near_clusters_count": len(near_clusters),
    }


# ============================================================
# EXPORT REPORT
# ============================================================

def export_report_txt(path, exact_groups, near_clusters, sim_reasons, summary):
    lines = []
    lines.append("=" * 70)
    lines.append("  AI SMART DUPLICATE FILE FINDER — REPORT")
    lines.append(f"  Generated: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    lines.append("=" * 70)
    lines.append("\nSUMMARY")
    lines.append(f"  Total files scanned   : {summary['total_files']}")
    lines.append(f"  Duplicate files found : {summary['duplicate_files']}")
    lines.append(f"  Wasted space          : {summary['wasted_str']}")
    lines.append(f"  Exact-match groups    : {summary['exact_groups_count']}")
    lines.append(f"  Near-duplicate groups : {summary['near_clusters_count']}")

    if exact_groups:
        lines.append("\n\nEXACT DUPLICATES (SHA-256)")
        lines.append("-" * 50)
        for i, grp in enumerate(exact_groups, 1):
            lines.append(f"\nGroup {i}:")
            for f in grp:
                lines.append(f"  {f}")

    if near_clusters:
        lines.append("\n\nNEAR-DUPLICATE CLUSTERS")
        lines.append("-" * 50)
        for i, cluster in enumerate(near_clusters, 1):
            keep, remove = recommend_cluster(cluster)
            lines.append(f"\nCluster {i}:")
            lines.append(f"  KEEP  : {keep[0]}")
            for f, _, _ in remove:
                lines.append(f"  REMOVE: {f}")
            for (f1, f2), reason in sim_reasons.items():
                if f1 in cluster and f2 in cluster:
                    lines.append(f"  [{reason}]  {os.path.basename(f1)} ↔ {os.path.basename(f2)}")

    with open(path, "w", encoding="utf-8") as fp:
        fp.write("\n".join(lines))


def export_report_csv(path, exact_groups, near_clusters, sim_reasons):
    rows = []
    for grp in exact_groups:
        for f in grp[1:]:
            meta = get_file_metadata(f)
            rows.append({
                "type": "exact", "cluster": "", "status": "REMOVE",
                "file": f, "size": meta["size_str"], "modified": meta["modified"],
                "similarity": "100%", "reason": "SHA-256 match",
            })

    for i, cluster in enumerate(near_clusters, 1):
        keep, remove = recommend_cluster(cluster)
        meta = get_file_metadata(keep[0])
        rows.append({
            "type": "near", "cluster": i, "status": "KEEP",
            "file": keep[0], "size": meta["size_str"], "modified": meta["modified"],
            "similarity": "", "reason": ", ".join(keep[2]),
        })
        for f, _, reasons in remove:
            meta     = get_file_metadata(f)
            pair_key = (keep[0], f) if (keep[0], f) in sim_reasons else (f, keep[0])
            reason   = sim_reasons.get(pair_key, "near-duplicate")
            rows.append({
                "type": "near", "cluster": i, "status": "REMOVE",
                "file": f, "size": meta["size_str"], "modified": meta["modified"],
                "similarity": reason, "reason": ", ".join(reasons),
            })

    with open(path, "w", newline="", encoding="utf-8") as fp:
        fieldnames = ["type", "cluster", "status", "file", "size", "modified", "similarity", "reason"]
        writer = csv.DictWriter(fp, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)


# ============================================================
# STANDARD MODE
# ============================================================

def analyze_folder(folder, threshold=0.65, progress_callback=None):
    def _log(msg):
        if progress_callback:
            progress_callback(msg)

    files = []
    for root, _, names in os.walk(folder):
        for n in names:
            files.append(os.path.join(root, n))

    _log(f"📂 Found {len(files)} files...")

    _log("🔐 Computing SHA-256 hashes...")
    size_groups = defaultdict(list)
    for f in files:
        try:
            size_groups[os.path.getsize(f)].append(f)
        except Exception:
            pass

    hash_groups = defaultdict(list)
    for group in size_groups.values():
        if len(group) < 2:
            continue
        for f in group:
            h = sha256_hash(f)
            if h:
                hash_groups[h].append(f)

    exact_duplicates = [g for g in hash_groups.values() if len(g) > 1]

    _log("📝 Extracting text and computing TF-IDF similarity...")
    text_files = [f for f in files if os.path.splitext(f)[1].lower() in TEXT_EXTENSIONS]
    texts, valid = [], []
    for f in text_files:
        txt = extract_text(f)
        print(f"[DEBUG] {os.path.basename(f)} → {len(txt)} chars")
        if len(txt) > 5:
            texts.append(txt)
            valid.append(f)

    _log(f"📝 {len(valid)} text files ready for comparison (threshold={threshold:.2f})...")

    near_pairs  = []
    sim_reasons = {}
    if len(texts) > 1:
        try:
            vec = TfidfVectorizer(stop_words=None, ngram_range=(1, 2), sublinear_tf=True, min_df=1)
            mat = cosine_similarity(vec.fit_transform(texts))
        except ValueError:
            vec = TfidfVectorizer(ngram_range=(1, 1), min_df=1)
            mat = cosine_similarity(vec.fit_transform(texts))
        for i in range(len(valid)):
            for j in range(i + 1, len(valid)):
                score = float(mat[i][j])
                print(f"[DEBUG] TF-IDF {os.path.basename(valid[i])} <-> {os.path.basename(valid[j])}: {score:.4f} (need>={threshold:.2f})")
                if score >= threshold:
                    near_pairs.append((valid[i], valid[j]))
                    sim_reasons[(valid[i], valid[j])] = f"TF-IDF similarity {int(score * 100)}%"

    # Include image pHash in Standard mode (was missing)
    _log("🖼️  Comparing images with perceptual hashing...")
    image_files = [f for f in files if os.path.splitext(f)[1].lower() in IMAGE_EXTENSIONS]
    hashes = {f: image_hash(f) for f in image_files}
    for i in range(len(image_files)):
        for j in range(i + 1, len(image_files)):
            h1 = hashes.get(image_files[i])
            h2 = hashes.get(image_files[j])
            if h1 and h2:
                score = 1 - (h1 - h2) / 64
                print(f"[DEBUG] pHash {os.path.basename(image_files[i])} <-> {os.path.basename(image_files[j])}: {score:.4f}")
                if score >= threshold:
                    near_pairs.append((image_files[i], image_files[j]))
                    sim_reasons[(image_files[i], image_files[j])] = f"Image pHash {int(score * 100)}%"

    clusters = cluster_pairs(near_pairs)
    summary  = compute_summary(exact_duplicates, clusters, files)
    return exact_duplicates, clusters, sim_reasons, summary, files


# ============================================================
# CROSS FORMAT MODE
# ============================================================

def analyze_cross_format(folder, threshold=0.65, progress_callback=None):
    def _log(msg):
        if progress_callback:
            progress_callback(msg)

    files = []
    for root, _, names in os.walk(folder):
        for n in names:
            files.append(os.path.join(root, n))

    _log(f"📂 Found {len(files)} files...")
    near_pairs  = []
    sim_reasons = {}

    _log("📝 Cross-format TF-IDF text extraction...")
    text_files = [f for f in files if os.path.splitext(f)[1].lower() in TEXT_EXTENSIONS]
    texts, valid = [], []
    for f in text_files:
        txt = extract_text(f)
        print(f"[DEBUG] {os.path.basename(f)} → {len(txt)} chars")
        if len(txt) > 5:
            texts.append(txt)
            valid.append(f)

    _log(f"📝 {len(valid)} text files ready (threshold={threshold:.2f})...")

    if len(texts) > 1:
        vec = TfidfVectorizer(stop_words=None, ngram_range=(1, 2), sublinear_tf=True, min_df=1)
        mat = cosine_similarity(vec.fit_transform(texts))
        for i in range(len(valid)):
            for j in range(i + 1, len(valid)):
                score = float(mat[i][j])
                print(f"[DEBUG] CrossFmt {os.path.basename(valid[i])} <-> {os.path.basename(valid[j])}: {score:.4f} (need>={threshold:.2f})")
                if score >= threshold:
                    near_pairs.append((valid[i], valid[j]))
                    sim_reasons[(valid[i], valid[j])] = f"Text similarity {int(score * 100)}%"

    _log("🖼️  Comparing images with perceptual hashing...")
    image_files = [f for f in files if os.path.splitext(f)[1].lower() in IMAGE_EXTENSIONS]
    hashes = {f: image_hash(f) for f in image_files}
    for i in range(len(image_files)):
        for j in range(i + 1, len(image_files)):
            h1 = hashes.get(image_files[i])
            h2 = hashes.get(image_files[j])
            if h1 and h2:
                score = 1 - (h1 - h2) / 64
                if score >= threshold:
                    near_pairs.append((image_files[i], image_files[j]))
                    sim_reasons[(image_files[i], image_files[j])] = f"Image pHash {int(score * 100)}%"

    _log("🎵 Fingerprinting audio files...")
    audio_files = [f for f in files if os.path.splitext(f)[1].lower() in AUDIO_EXTENSIONS]
    audio_fps   = {}
    for f in audio_files:
        fp = audio_fingerprint(f)
        if fp is not None:
            audio_fps[f] = fp

    af_list = list(audio_fps.keys())
    for i in range(len(af_list)):
        for j in range(i + 1, len(af_list)):
            score = audio_similarity(audio_fps[af_list[i]], audio_fps[af_list[j]])
            if score >= threshold:
                near_pairs.append((af_list[i], af_list[j]))
                sim_reasons[(af_list[i], af_list[j])] = f"Audio fingerprint {int(score * 100)}%"

    clusters = cluster_pairs(near_pairs)
    summary  = compute_summary([], clusters, files)
    return clusters, sim_reasons, summary, files


# ============================================================
# DEEP AI MODE
# ============================================================

def analyze_deep_ai(folder, threshold=0.6, progress_callback=None):
    def _log(msg):
        if progress_callback:
            progress_callback(msg)

    files = []
    for root, _, names in os.walk(folder):
        for n in names:
            files.append(os.path.join(root, n))

    _log(f"📂 Found {len(files)} files...")
    near_pairs  = {}
    sim_reasons = {}

    # DOCUMENT SEMANTIC EMBEDDING
    text_files = [f for f in files if os.path.splitext(f)[1].lower() in TEXT_EXTENSIONS]
    texts, valid = [], []
    for f in text_files:
        txt = extract_text(f)
        print(f"[DEBUG] {os.path.basename(f)} → {len(txt)} chars")
        if len(txt) >= 5:
            texts.append(txt)
            valid.append(f)

    _log(f"📝 {len(valid)} text files ready for semantic embedding (threshold={threshold:.2f})...")

    if len(valid) > 1:
        _log("🧠 Loading Sentence Transformer model (all-MiniLM-L6-v2)...")
        from sentence_transformers import SentenceTransformer
        import numpy as np
        model = SentenceTransformer("all-MiniLM-L6-v2")
        _log("⚙️  Encoding documents into vector embeddings...")
        emb   = model.encode(texts, convert_to_numpy=True)
        norms = np.linalg.norm(emb, axis=1, keepdims=True)
        emb   = emb / (norms + 1e-10)
        index = build_faiss_index(emb)
        k = min(10, len(valid))
        distances, indices = index.search(emb.astype("float32"), k)
        _log("🔗 Finding semantic duplicate pairs...")
        for i in range(len(valid)):
            for rank in range(1, k):
                j     = int(indices[i][rank])
                score = float(distances[i][rank])
                if j <= i:
                    continue
                if score >= threshold:
                    f1, f2 = valid[i], valid[j]
                    near_pairs[(f1, f2)]  = score
                    sim_reasons[(f1, f2)] = f"Semantic similarity {int(score * 100)}%"

    # IMAGE AI — CLIP + pHash
    image_files = [f for f in files if os.path.splitext(f)[1].lower() in IMAGE_EXTENSIONS]
    try:
        import torch
        import clip as openai_clip
        import numpy as np
        _log("🎨 Loading CLIP vision model (ViT-B/32)...")
        device = "cuda" if torch.cuda.is_available() else "cpu"
        clip_model, preprocess = openai_clip.load("ViT-B/32", device=device)
        embeddings = []
        valid_imgs = []
        _log("🖼️  Encoding images with CLIP neural vision...")
        for f in image_files:
            try:
                img = preprocess(Image.open(f)).unsqueeze(0).to(device)
                with torch.no_grad():
                    e = clip_model.encode_image(img)
                    e = e / e.norm(dim=-1, keepdim=True)
                embeddings.append(e.cpu().numpy()[0])
                valid_imgs.append(f)
            except Exception:
                pass

        if embeddings:
            embeddings = np.array(embeddings)
            sim        = embeddings @ embeddings.T
            for i in range(len(valid_imgs)):
                for j in range(i + 1, len(valid_imgs)):
                    clip_score = float(sim[i][j])
                    if clip_score >= threshold:
                        f1, f2      = valid_imgs[i], valid_imgs[j]
                        h1, h2      = image_hash(f1), image_hash(f2)
                        if h1 and h2:
                            phash_score = 1 - ((h1 - h2) / 64)
                            if phash_score >= (threshold * 0.85):
                                near_pairs[(f1, f2)]  = phash_score
                                sim_reasons[(f1, f2)] = f"CLIP+pHash {int(phash_score * 100)}%"
    except Exception as e:
        _log(f"⚠️  CLIP not available ({e}), skipping neural image scan.")

    # AUDIO
    _log("🎵 Fingerprinting audio files...")
    audio_files = [f for f in files if os.path.splitext(f)[1].lower() in AUDIO_EXTENSIONS]
    audio_fps   = {}
    for f in audio_files:
        fp = audio_fingerprint(f)
        if fp is not None:
            audio_fps[f] = fp

    af_list = list(audio_fps.keys())
    for i in range(len(af_list)):
        for j in range(i + 1, len(af_list)):
            score = audio_similarity(audio_fps[af_list[i]], audio_fps[af_list[j]])
            if score >= threshold:
                f1, f2 = af_list[i], af_list[j]
                near_pairs[(f1, f2)]  = score
                sim_reasons[(f1, f2)] = f"Audio MFCC {int(score * 100)}%"

    clusters = cluster_pairs(list(near_pairs.keys()))
    summary  = compute_summary([], clusters, files)
    return clusters, near_pairs, sim_reasons, summary, files


# ============================================================
# CLUSTER EXPLANATION
# ============================================================

def generate_cluster_explanation(cluster, keep_file, remove_files, similarity_reasons):
    lines = []
    for (f1, f2), reason in similarity_reasons.items():
        if f1 in cluster and f2 in cluster:
            lines.append(f"{os.path.basename(f1)} ↔ {os.path.basename(f2)} — {reason}")
    similarity_text = "\n".join(lines)
    keep_line       = f"📌 Keep: {os.path.basename(keep_file[0])}"
    return similarity_text, keep_line